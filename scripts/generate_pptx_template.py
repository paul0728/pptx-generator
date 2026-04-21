"""PPTX Generator — reference implementation for the pptx-generator skill.

Pipeline (see SKILL.md):
    Phase 1  load & validate slides.json
    Phase 2  render Mermaid diagrams (parallel, cached, CJK fallback)
    Phase 3  assemble PPTX (name-based layout lookup, CJK-aware auto-fit)
    Phase 4  quality verification

Usage:
    python generate_pptx_template.py --json slides.json \\
        [--template template.pptx] [--out output.pptx] [-v]

Dependencies:
    pip install python-pptx requests Pillow
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import logging
import re
import sys
import tempfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from contextlib import suppress
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

import json
import requests
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger("pptx_generator")

# ── Constants ─────────────────────────────────────────────────────────
SKILL_DIR = Path(__file__).resolve().parent.parent  # repo root
DEFAULT_TEMPLATE = SKILL_DIR / "assets" / "default-template.pptx"
MERMAID_CACHE_DIR = SKILL_DIR / ".cache" / "mermaid"
MERMAID_URL = "https://mermaid.ink/img/{encoded}?width=1600"
MERMAID_TIMEOUT = 30
MERMAID_MAX_WORKERS = 4
MERMAID_RETRIES = 3

EMU_PER_INCH = 914400
MERMAID_MAX_NODES = 15  # warn above this density

DEFAULT_BRAND_COLOR = "#2B579A"
DEFAULT_FONT_LATIN = "Calibri"
DEFAULT_FONT_CJK = "微軟正黑體"
DEFAULT_FONT_CODE = "Consolas"

TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
BODY_TYPES = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
SUBTITLE_TYPES = {PP_PLACEHOLDER.SUBTITLE}

# Layout name candidates (exact match first, then substring, then fallback index).
LAYOUT_HINTS: dict[str, tuple[tuple[str, ...], int]] = {
    # role : (name candidates, fallback index)
    # Order matters: first name that matches wins (exact > substring).
    # Prefer simpler cover layouts over photo-heavy variants (that have empty
    # picture placeholders which render as white boxes over the logo).
    "title": (("Title Slide", "新版封面", "標題投影片", "封面"), 0),
    "bullet": (("Title and Content", "標題及內容", "一般", "內容"), 1),
    "section": (
        ("Section Header", "章節標題", "章節", "單元封面", "單元封面(右)", "單元封面-綠底"),
        1,
    ),
    "diagram": (("Title and Content", "標題及內容", "一般"), 1),
    "code": (("Title and Content", "標題及內容", "一般"), 1),
    "two_col": (("Two Content", "兩欄", "雙欄", "Title and Content", "一般"), 1),
    "table": (("Title and Content", "標題及內容", "一般"), 1),
    "image": (("Title and Content", "標題及內容", "一般"), 1),
    "kpi": (("Title and Content", "標題及內容", "一般"), 1),
    "ending": (("Title Slide", "結尾", "Ending", "封面"), 0),
}


# ── Brand / chrome config ─────────────────────────────────────────────
@dataclass
class BrandConfig:
    color: str = DEFAULT_BRAND_COLOR  # hex #RRGGBB
    font: str | None = None  # overrides default body/title font
    footer: str | None = None  # bottom-left footer text
    version: str | None = None  # bottom-right version/date
    watermark: str | None = None  # diagonal watermark
    page_numbers: bool = False  # show page x/total
    skip_chrome_on_title: bool = True  # don't add footer on cover/ending

    @staticmethod
    def _parse_hex(value: str) -> RGBColor:
        v = value.lstrip("#")
        if len(v) != 6:
            raise ValueError(f"brand color must be #RRGGBB, got {value!r}")
        return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))

    @property
    def rgb(self) -> RGBColor:
        return self._parse_hex(self.color)


# ── Data model ────────────────────────────────────────────────────────
@dataclass
class Slide:
    id: int
    type: str
    content: dict
    _errors: list[str] = field(default_factory=list)

    VALID_TYPES = {
        "title_slide",
        "outline_slide",
        "section_slide",
        "bullet_points",
        "architecture_diagram",
        "code_demo",
        "ending_slide",
        "two_column",
        "table",
        "image_slide",
        "kpi_slide",
    }

    @classmethod
    def parse(cls, raw: dict) -> "Slide":
        sid = raw.get("id")
        stype = raw.get("type")
        content = raw.get("content") or {}
        if not isinstance(sid, int):
            raise ValueError(f"slide.id must be int, got {sid!r}")
        if stype not in cls.VALID_TYPES:
            raise ValueError(f"slide {sid}: unknown type {stype!r}")
        if not isinstance(content, dict):
            raise ValueError(f"slide {sid}: content must be object")
        return cls(id=sid, type=stype, content=content)


def load_slides(path: Path) -> tuple[list[Slide], dict]:
    data = json.loads(path.read_text(encoding="utf-8"))
    if "slides" not in data or not isinstance(data["slides"], list):
        raise ValueError(f"{path}: missing 'slides' array")
    meta = data.get("presentation_metadata") or {}
    return [Slide.parse(s) for s in data["slides"]], meta


# ── Phase 2: Mermaid rendering ────────────────────────────────────────
_CJK_RE = re.compile(r"[\u3000-\u9fff\uac00-\ud7af\uff00-\uffef]")


def _has_cjk(text: str) -> bool:
    return bool(_CJK_RE.search(text))


def _strip_cjk_labels(mermaid: str) -> str:
    """Replace CJK-only node labels with ASCII placeholders as a last-resort fallback."""

    def repl(m: re.Match) -> str:
        label = m.group(1)
        ascii_only = re.sub(r"[\u3000-\u9fff\uac00-\ud7af\uff00-\uffef]+", "Node", label).strip()
        return f"[{ascii_only or 'Node'}]"

    return re.sub(r"\[([^\]]+)\]", repl, mermaid)


def _count_mermaid_nodes(code: str) -> int:
    """Rough node count: unique tokens that start a node shape bracket."""
    nodes = set(re.findall(r"(?:^|\s|-->|---|==>|-\.->)\s*([A-Za-z0-9_]+)\s*[\[\(\{]", code))
    # Also count bare id references on the right of -->
    nodes |= set(re.findall(r"-->\s*([A-Za-z0-9_]+)\b", code))
    return len(nodes)


def _mermaid_request(mermaid_code: str) -> bytes:
    encoded = base64.urlsafe_b64encode(mermaid_code.encode("utf-8")).decode("ascii")
    url = MERMAID_URL.format(encoded=encoded)
    last_exc: Exception | None = None
    for attempt in range(1, MERMAID_RETRIES + 1):
        try:
            r = requests.get(url, timeout=MERMAID_TIMEOUT)
            r.raise_for_status()
            if len(r.content) < 200:
                raise RuntimeError(f"suspiciously small response ({len(r.content)} bytes)")
            return r.content
        except Exception as exc:  # noqa: BLE001 — retry regardless
            last_exc = exc
            logger.debug("mermaid attempt %d/%d failed: %s", attempt, MERMAID_RETRIES, exc)
    assert last_exc is not None
    raise last_exc


def render_mermaid(mermaid_code: str, out_path: Path) -> bool:
    """Render Mermaid → PNG with disk cache and CJK fallback.

    Returns True on success. The output file is guaranteed to exist on success.
    """
    out_path = Path(out_path)
    MERMAID_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    cache_key = hashlib.sha256(mermaid_code.encode("utf-8")).hexdigest()
    cached = MERMAID_CACHE_DIR / f"{cache_key}.png"

    if cached.exists():
        out_path.write_bytes(cached.read_bytes())
        logger.debug("mermaid cache hit: %s", cache_key[:8])
        return True

    candidates: list[str] = [mermaid_code]
    if _has_cjk(mermaid_code):
        candidates.append(_strip_cjk_labels(mermaid_code))

    for i, code in enumerate(candidates):
        try:
            data = _mermaid_request(code)
            out_path.write_bytes(data)
            cached.write_bytes(data)
            if i > 0:
                logger.warning("mermaid CJK fallback used (ASCII labels substituted)")
            return True
        except Exception as exc:  # noqa: BLE001
            logger.warning("mermaid render attempt failed: %s", exc)

    return False


def render_all_diagrams(slides: list[Slide], tmp_dir: Path) -> dict[int, Path]:
    """Render every slide's mermaid_code in parallel. Returns {slide_id: png_path}."""
    jobs: dict[int, Path] = {}
    for s in slides:
        mc = s.content.get("mermaid_code")
        if mc:
            n = _count_mermaid_nodes(mc)
            if n > MERMAID_MAX_NODES:
                logger.warning(
                    "slide %d: mermaid has ~%d nodes (> %d); diagram may be cramped — consider splitting",
                    s.id,
                    n,
                    MERMAID_MAX_NODES,
                )
            jobs[s.id] = tmp_dir / f"diagram_{s.id}.png"

    if not jobs:
        return {}

    success: dict[int, Path] = {}
    with ThreadPoolExecutor(max_workers=min(MERMAID_MAX_WORKERS, len(jobs))) as pool:
        futures = {
            pool.submit(
                render_mermaid, next(s.content["mermaid_code"] for s in slides if s.id == sid), path
            ): sid
            for sid, path in jobs.items()
        }
        for fut in as_completed(futures):
            sid = futures[fut]
            try:
                if fut.result():
                    success[sid] = jobs[sid]
                else:
                    logger.error("slide %d: diagram render failed", sid)
            except Exception as exc:  # noqa: BLE001
                logger.error("slide %d: diagram render crashed: %s", sid, exc)

    logger.info("mermaid: %d/%d succeeded", len(success), len(jobs))
    return success


# ── Phase 3-A: Layout & placeholder lookup ────────────────────────────
def find_layout(prs: Presentation, role: str):
    names, fallback_idx = LAYOUT_HINTS[role]
    # 1. exact name match
    for name in names:
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
    # 2. substring match (case-insensitive)
    for name in names:
        for layout in prs.slide_layouts:
            if name.lower() in layout.name.lower():
                return layout
    # 3. fallback
    idx = min(fallback_idx, len(prs.slide_layouts) - 1)
    logger.warning(
        "layout for role %r not found in template; falling back to index %d (%s)",
        role,
        idx,
        prs.slide_layouts[idx].name,
    )
    return prs.slide_layouts[idx]


def find_placeholder(slide, *, kinds: Iterable, allow_any_body: bool = False):
    """Return first placeholder matching one of the given PP_PLACEHOLDER kinds."""
    kind_set = set(kinds)
    for ph in slide.placeholders:
        if ph.placeholder_format.type in kind_set:
            return ph
    if allow_any_body:
        # Fallback: any non-title placeholder.
        for ph in slide.placeholders:
            if ph.placeholder_format.type not in TITLE_TYPES:
                return ph
    return None


# ── Phase 3-B: CJK-aware auto-fit ─────────────────────────────────────
def _visual_width_in(text: str, pt: int) -> float:
    """Estimate visual width (inches) of a single line at given point size."""
    char_w = pt / 72.0  # 1 CJK char ≈ 1em
    total = 0.0
    for c in text:
        if c == "\t":
            total += char_w * 2.2
        elif ord(c) > 127:
            total += char_w
        else:
            total += char_w * 0.55
    return total


def auto_fit_text(shape, text: str, *, max_pt: int, min_pt: int, is_code: bool = False) -> None:
    """Insert text and shrink font until it fits. CJK-aware; forces Consolas for code."""
    shape.text = text
    if not text.strip():
        return

    tf = shape.text_frame
    tf.word_wrap = True
    ph_w = shape.width / EMU_PER_INCH
    ph_h = shape.height / EMU_PER_INCH

    chosen_pt = min_pt
    for pt in range(max_pt, min_pt - 1, -1):
        total_lines = 0
        for line in text.split("\n"):
            w = _visual_width_in(line, pt)
            total_lines += max(1, int(w / ph_w + 0.999))
        if total_lines * (pt / 72.0 * 1.35) <= ph_h:
            chosen_pt = pt
            break
    else:
        logger.warning(
            "shape overflow: content at min %dpt still does not fit (consider splitting slide)",
            min_pt,
        )

    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(chosen_pt)
            if is_code:
                r.font.name = "Consolas"


def _set_placeholder_text(
    slide,
    *,
    kinds: Iterable,
    text: str,
    max_pt: int,
    min_pt: int,
    is_code: bool = False,
    allow_any_body: bool = False,
) -> bool:
    ph = find_placeholder(slide, kinds=kinds, allow_any_body=allow_any_body)
    if ph is None:
        logger.warning("placeholder not found (kinds=%s); text dropped: %r", kinds, text[:30])
        return False
    auto_fit_text(ph, text, max_pt=max_pt, min_pt=min_pt, is_code=is_code)
    return True


def set_title(slide, text: str) -> None:
    if not text:
        return
    _set_placeholder_text(
        slide,
        kinds=TITLE_TYPES,
        text=text,
        max_pt=24,
        min_pt=12,
        allow_any_body=True,
    )


def set_subtitle(slide, text: str) -> None:
    if not text:
        return
    _set_placeholder_text(
        slide,
        kinds=SUBTITLE_TYPES,
        text=text,
        max_pt=14,
        min_pt=8,
        allow_any_body=True,
    )


# ── Bullet rendering with hierarchy ───────────────────────────────────
def _parse_bullet(line: str) -> tuple[int, str]:
    """Detect indent level (0-4) from leading spaces (2 spaces = 1 level)."""
    stripped = line.lstrip(" ")
    indent = len(line) - len(stripped)
    # Also accept leading "- " or "• " after indent; strip marker.
    stripped = re.sub(r"^([-•·*]|\d+[.)])\s+", "", stripped)
    return min(indent // 2, 4), stripped


def set_bullets(slide, points: list[str], *, max_pt: int = 18, min_pt: int = 7) -> None:
    ph = find_placeholder(slide, kinds=BODY_TYPES, allow_any_body=True)
    if ph is None:
        logger.warning("body placeholder not found; bullets dropped")
        return
    tf = ph.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, raw in enumerate(points):
        level, text = _parse_bullet(raw)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level

    # Reuse auto_fit_text's sizing logic by synthesizing a measured block.
    flat = "\n".join(_parse_bullet(p)[1] for p in points)
    ph_w = ph.width / EMU_PER_INCH
    ph_h = ph.height / EMU_PER_INCH
    chosen_pt = min_pt
    for pt in range(max_pt, min_pt - 1, -1):
        total = 0
        for line in flat.split("\n"):
            w = _visual_width_in(line, pt)
            total += max(1, int(w / ph_w + 0.999))
        if total * (pt / 72.0 * 1.35) <= ph_h:
            chosen_pt = pt
            break
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(chosen_pt)


# ── Template loading ──────────────────────────────────────────────────
def load_template(template_path: Path | None) -> Presentation:
    """Load template, prefer explicit → default asset → blank."""
    if template_path and template_path.exists():
        logger.info("using template: %s", template_path)
        prs = Presentation(str(template_path))
    elif DEFAULT_TEMPLATE.exists():
        logger.info("using default template: %s", DEFAULT_TEMPLATE)
        prs = Presentation(str(DEFAULT_TEMPLATE))
    else:
        logger.info("no template; using blank presentation")
        prs = Presentation()

    # Remove any pre-existing slides (keep layouts & master).
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        rel_id = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        ) or sld_id.get("r:id")
        if rel_id is not None:
            with suppress(KeyError):
                prs.part.drop_rel(rel_id)
        xml_slides.remove(sld_id)
    return prs


# ── Slide builders ────────────────────────────────────────────────────
def build_title_slide(prs: Presentation, s: Slide) -> None:
    slide = prs.slides.add_slide(find_layout(prs, "title"))
    set_title(slide, s.content.get("title", ""))
    set_subtitle(slide, s.content.get("sub_title", ""))
    _apply_notes(slide, s.content.get("notes"))


def build_outline_or_bullets(prs: Presentation, s: Slide) -> None:
    slide = prs.slides.add_slide(find_layout(prs, "bullet"))
    set_title(slide, s.content.get("title", ""))
    points = s.content.get("points", []) or []
    if points:
        set_bullets(slide, points)
    _apply_notes(slide, s.content.get("notes"))


def build_section_slide(prs: Presentation, s: Slide) -> None:
    slide = prs.slides.add_slide(find_layout(prs, "section"))
    set_title(slide, s.content.get("title", ""))
    set_subtitle(slide, s.content.get("sub_title", ""))
    _apply_notes(slide, s.content.get("notes"))


def build_diagram_slide(prs: Presentation, s: Slide, diagrams: dict[int, Path]) -> None:
    slide = prs.slides.add_slide(find_layout(prs, "diagram"))
    set_title(slide, s.content.get("title", ""))

    img_path = diagrams.get(s.id)
    slide_w_in = prs.slide_width / EMU_PER_INCH

    if img_path and img_path.exists():
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        # Work in inches directly using aspect ratio; ignore DPI.
        aspect = iw / ih
        max_w, max_h = slide_w_in - 1.0, 3.7  # leave 0.5" margins L/R, room for desc
        if aspect >= max_w / max_h:
            fw, fh = max_w, max_w / aspect
        else:
            fw, fh = max_h * aspect, max_h
        left = (slide_w_in - fw) / 2
        slide.shapes.add_picture(
            str(img_path),
            Inches(left),
            Inches(1.1),
            width=Inches(fw),
            height=Inches(fh),
        )
        _set_picture_alt(slide.shapes[-1], s.content.get("title", ""))

        # Description goes in a textbox BELOW the image — not in body placeholder
        # (avoids overlap since body placeholder area is reused by the image).
        desc = s.content.get("description", "")
        if desc:
            desc_top = 1.1 + fh + 0.15
            desc_h = max(0.4, 6.9 - desc_top)
            tb = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(desc_top),
                Inches(slide_w_in - 1.0),
                Inches(desc_h),
            )
            auto_fit_text(tb, desc, max_pt=12, min_pt=8)
    else:
        _set_placeholder_text(
            slide,
            kinds=BODY_TYPES,
            text=f"[Diagram rendering failed]\n\n{s.content.get('description', '')}",
            max_pt=14,
            min_pt=8,
            allow_any_body=True,
        )
    _apply_notes(slide, s.content.get("notes"))


def build_code_slide(prs: Presentation, s: Slide) -> None:
    slide = prs.slides.add_slide(find_layout(prs, "code"))
    set_title(slide, s.content.get("title", ""))
    _set_placeholder_text(
        slide,
        kinds=BODY_TYPES,
        text=s.content.get("code", ""),
        max_pt=12,
        min_pt=6,
        is_code=True,
        allow_any_body=True,
    )
    _apply_notes(slide, s.content.get("notes"))


def build_ending_slide(prs: Presentation, s: Slide) -> None:
    slide = prs.slides.add_slide(find_layout(prs, "ending"))
    set_title(slide, s.content.get("title", "Thank You"))
    set_subtitle(slide, s.content.get("sub_title", ""))
    _apply_notes(slide, s.content.get("notes"))


# ── New slide types: two_column / table / image_slide / kpi_slide ────
def build_two_column_slide(prs: Presentation, s: Slide) -> None:
    """Left/right columns with heading + bullet points."""
    slide = prs.slides.add_slide(find_layout(prs, "two_col"))
    set_title(slide, s.content.get("title", ""))

    left = s.content.get("left") or {}
    right = s.content.get("right") or {}
    sw = prs.slide_width / EMU_PER_INCH
    sh = prs.slide_height / EMU_PER_INCH
    margin = 0.5
    gap = 0.3
    top = 1.5
    col_w = (sw - margin * 2 - gap) / 2
    col_h = sh - top - 0.8

    for i, col in enumerate((left, right)):
        left_in = margin + i * (col_w + gap)
        # Column heading
        if col.get("heading"):
            hb = slide.shapes.add_textbox(
                Inches(left_in),
                Inches(top),
                Inches(col_w),
                Inches(0.5),
            )
            auto_fit_text(hb, col["heading"], max_pt=16, min_pt=10)
            for p in hb.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
            body_top = top + 0.55
            body_h = col_h - 0.55
        else:
            body_top = top
            body_h = col_h
        # Column body
        tb = slide.shapes.add_textbox(
            Inches(left_in),
            Inches(body_top),
            Inches(col_w),
            Inches(body_h),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        points = col.get("points") or ([col["text"]] if col.get("text") else [])
        if not points:
            continue
        tf.clear()
        for j, raw in enumerate(points):
            level, text = _parse_bullet(raw)
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = text
            p.level = level
        # size pass
        flat = "\n".join(_parse_bullet(p)[1] for p in points)
        chosen_pt = 7
        for pt in range(16, 6, -1):
            total = 0
            for line in flat.split("\n"):
                w = _visual_width_in(line, pt)
                total += max(1, int(w / col_w + 0.999))
            if total * (pt / 72.0 * 1.35) <= body_h:
                chosen_pt = pt
                break
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(chosen_pt)
    _apply_notes(slide, s.content.get("notes"))


def build_table_slide(prs: Presentation, s: Slide) -> None:
    """Render a data table. content: headers: [...], rows: [[...], ...]."""
    slide = prs.slides.add_slide(find_layout(prs, "table"))
    set_title(slide, s.content.get("title", ""))

    headers = s.content.get("headers") or []
    rows = s.content.get("rows") or []
    if not headers or not rows:
        logger.warning("slide %d: table missing headers or rows", s.id)
        _apply_notes(slide, s.content.get("notes"))
        return

    sw = prs.slide_width / EMU_PER_INCH
    margin = 0.5
    top, width = 1.5, sw - margin * 2
    max_h = (prs.slide_height / EMU_PER_INCH) - top - 0.8
    row_h = min(0.45, max_h / (len(rows) + 1))

    table_shape = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(headers),
        left=Inches(margin),
        top=Inches(top),
        width=Inches(width),
        height=Inches(row_h * (len(rows) + 1)),
    )
    tbl = table_shape.table
    # header row
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(12)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2B, 0x57, 0x9A)
    # data rows
    body_pt = 11 if len(rows) <= 8 else 10 if len(rows) <= 12 else 9
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row[: len(headers)]):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = "" if val is None else str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(body_pt)
    _apply_notes(slide, s.content.get("notes"))


def build_image_slide(prs: Presentation, s: Slide) -> None:
    """Display an image with optional caption. content: image_path, caption."""
    slide = prs.slides.add_slide(find_layout(prs, "image"))
    set_title(slide, s.content.get("title", ""))

    img_path = s.content.get("image_path") or s.content.get("image")
    caption = s.content.get("caption", "")
    sw = prs.slide_width / EMU_PER_INCH

    if not img_path or not Path(img_path).exists():
        _set_placeholder_text(
            slide,
            kinds=BODY_TYPES,
            text=f"[Image not found: {img_path}]\n{caption}",
            max_pt=14,
            min_pt=8,
            allow_any_body=True,
        )
        _apply_notes(slide, s.content.get("notes"))
        return

    with PILImage.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    max_w, max_h = sw - 1.0, 4.5
    if aspect >= max_w / max_h:
        fw, fh = max_w, max_w / aspect
    else:
        fw, fh = max_h * aspect, max_h
    left = (sw - fw) / 2
    pic = slide.shapes.add_picture(
        str(img_path),
        Inches(left),
        Inches(1.1),
        width=Inches(fw),
        height=Inches(fh),
    )
    _set_picture_alt(pic, caption or s.content.get("title", "image"))

    if caption:
        cap_top = 1.1 + fh + 0.15
        tb = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(cap_top),
            Inches(sw - 1.0),
            Inches(0.5),
        )
        auto_fit_text(tb, caption, max_pt=12, min_pt=8)
        for p in tb.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.italic = True
    _apply_notes(slide, s.content.get("notes"))


def build_kpi_slide(prs: Presentation, s: Slide, brand: BrandConfig) -> None:
    """Grid of KPI cards. content: kpis: [{label, value, delta?, unit?}]."""
    slide = prs.slides.add_slide(find_layout(prs, "kpi"))
    set_title(slide, s.content.get("title", ""))

    kpis = s.content.get("kpis") or []
    if not kpis:
        _apply_notes(slide, s.content.get("notes"))
        return

    n = min(len(kpis), 6)
    cols = 2 if n <= 2 else 3
    rows = (n + cols - 1) // cols
    sw = prs.slide_width / EMU_PER_INCH
    sh = prs.slide_height / EMU_PER_INCH
    margin, gap = 0.6, 0.3
    top = 1.5
    card_w = (sw - margin * 2 - gap * (cols - 1)) / cols
    avail_h = sh - top - 0.8
    card_h = min(1.8, (avail_h - gap * (rows - 1)) / rows)

    for i, k in enumerate(kpis[:n]):
        r, c = divmod(i, cols)
        x = margin + c * (card_w + gap)
        y = top + r * (card_h + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(card_w),
            Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
        card.line.color.rgb = brand.rgb
        card.line.width = Emu(12700)  # 1pt
        card.text_frame.text = ""

        # value (large)
        val = f"{k.get('value', '')}{k.get('unit', '')}"
        val_box = slide.shapes.add_textbox(
            Inches(x),
            Inches(y + 0.15),
            Inches(card_w),
            Inches(card_h * 0.5),
        )
        val_box.text_frame.text = val
        for p in val_box.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(36 if card_w > 2.2 else 28)
                run.font.bold = True
                run.font.color.rgb = brand.rgb

        # label
        lbl_box = slide.shapes.add_textbox(
            Inches(x),
            Inches(y + card_h * 0.55),
            Inches(card_w),
            Inches(card_h * 0.25),
        )
        lbl_box.text_frame.text = str(k.get("label", ""))
        for p in lbl_box.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        # delta (optional)
        delta = k.get("delta")
        if delta:
            d_box = slide.shapes.add_textbox(
                Inches(x),
                Inches(y + card_h * 0.78),
                Inches(card_w),
                Inches(card_h * 0.2),
            )
            d_box.text_frame.text = str(delta)
            is_down = str(delta).strip().startswith(("-", "▼", "↓"))
            color = RGBColor(0xC0, 0x39, 0x2B) if is_down else RGBColor(0x2E, 0x8B, 0x57)
            for p in d_box.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = color
    _apply_notes(slide, s.content.get("notes"))


# ── Chrome: footer / page numbers / watermark ─────────────────────────
def _is_cover_slide(s: Slide) -> bool:
    return s.type in ("title_slide", "ending_slide")


def apply_chrome(prs: Presentation, slides: list[Slide], brand: BrandConfig) -> None:
    """Add footer text, page numbers, and watermark to body slides."""
    sw = prs.slide_width / EMU_PER_INCH
    sh = prs.slide_height / EMU_PER_INCH
    total = len(prs.slides)

    for idx, (slide, spec) in enumerate(zip(prs.slides, slides), start=1):
        if brand.skip_chrome_on_title and _is_cover_slide(spec):
            continue

        # watermark (diagonal, very light)
        if brand.watermark:
            wm = slide.shapes.add_textbox(
                Inches(sw * 0.15),
                Inches(sh * 0.35),
                Inches(sw * 0.7),
                Inches(sh * 0.3),
            )
            wm.text_frame.text = brand.watermark
            for p in wm.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(72)
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
            # rotate -30°
            with suppress(Exception):
                wm.rotation = -30.0
            # push to back
            with suppress(Exception):
                spTree = wm._element.getparent()
                spTree.remove(wm._element)
                spTree.insert(2, wm._element)

        # footer (bottom-left)
        if brand.footer:
            fb = slide.shapes.add_textbox(
                Inches(0.4),
                Inches(sh - 0.35),
                Inches(sw * 0.5),
                Inches(0.3),
            )
            fb.text_frame.text = brand.footer
            for p in fb.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
                    r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

        # version / date (bottom-center)
        if brand.version:
            vb = slide.shapes.add_textbox(
                Inches(sw * 0.3),
                Inches(sh - 0.35),
                Inches(sw * 0.4),
                Inches(0.3),
            )
            vb.text_frame.text = brand.version
            for p in vb.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(9)
                    r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

        # page numbers (bottom-right)
        if brand.page_numbers:
            pb = slide.shapes.add_textbox(
                Inches(sw - 1.2),
                Inches(sh - 0.35),
                Inches(1.0),
                Inches(0.3),
            )
            pb.text_frame.text = f"{idx} / {total}"
            for p in pb.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT
                for r in p.runs:
                    r.font.size = Pt(9)
                    r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)


def apply_brand_font(prs: Presentation, brand: BrandConfig) -> None:
    """Apply brand font to all non-code text runs across every slide."""
    if not brand.font:
        return
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name == DEFAULT_FONT_CODE:
                        continue
                    r.font.name = brand.font


BUILDERS = {
    "title_slide": build_title_slide,
    "outline_slide": build_outline_or_bullets,
    "bullet_points": build_outline_or_bullets,
    "section_slide": build_section_slide,
    "code_demo": build_code_slide,
    "ending_slide": build_ending_slide,
    "two_column": build_two_column_slide,
    "table": build_table_slide,
    "image_slide": build_image_slide,
    # kpi_slide receives brand and is dispatched specially in generate()
}


# ── Helpers ───────────────────────────────────────────────────────────
def _apply_notes(slide, notes: str | None) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def strip_unused_picture_placeholders(slide) -> None:
    """Remove empty picture placeholders inherited from the layout.

    Picture placeholders render as opaque white boxes when no image is set,
    which can hide template graphics (e.g. the corporate logo) positioned
    beneath them. Call this after populating a slide, before saving.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER as _PP

    to_remove = []
    for ph in list(slide.placeholders):
        fmt = ph.placeholder_format
        if fmt.type != _PP.PICTURE:
            continue
        # If there's no image blip inside, treat as empty.
        has_image = bool(
            ph._element.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip")
        )
        if not has_image:
            to_remove.append(ph._element)
    for el in to_remove:
        el.getparent().remove(el)


def _set_picture_alt(picture, alt_text: str) -> None:
    """Set accessibility alt-text on a picture shape."""
    with suppress(Exception):
        nv = picture._element.nvPicPr.cNvPr  # type: ignore[attr-defined]
        nv.set("descr", alt_text)
        nv.set("title", alt_text)


def save_with_retry(prs: Presentation, out: Path, *, max_attempts: int = 3) -> Path:
    """Save with retry; adds -v2, -v3 suffix when the file is locked."""
    out.parent.mkdir(parents=True, exist_ok=True)
    last_exc: Exception | None = None
    for attempt in range(1, max_attempts + 1):
        target = out if attempt == 1 else out.with_stem(f"{out.stem}-v{attempt}")
        try:
            prs.save(str(target))
            return target
        except PermissionError as exc:
            last_exc = exc
            logger.warning("save locked (%s); retrying as %s", out.name, target.name)
    assert last_exc is not None
    raise last_exc


# ── Phase 4: verification ─────────────────────────────────────────────
def verify_output(path: Path, expected_slides: int) -> list[str]:
    issues: list[str] = []
    prs = Presentation(str(path))
    if len(prs.slides) != expected_slides:
        issues.append(f"expected {expected_slides} slides, got {len(prs.slides)}")
    for i, slide in enumerate(prs.slides, start=1):
        has_text = any(
            getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
            for shape in slide.shapes
        )
        has_picture = any(
            shape.shape_type == 13 for shape in slide.shapes
        )  # MSO_SHAPE_TYPE.PICTURE
        if not (has_text or has_picture):
            issues.append(f"slide {i}: no visible content")
    return issues


# ── Orchestration ─────────────────────────────────────────────────────
def generate(
    json_path: Path,
    template_path: Path | None,
    output_path: Path,
    brand: BrandConfig | None = None,
) -> Path:
    logger.info("Phase 1: parsing %s", json_path)
    slides, meta = load_slides(json_path)
    logger.info("  %d slides", len(slides))

    brand = brand or BrandConfig()
    # metadata can supply footer/version when not set via CLI
    if not brand.footer and meta.get("title"):
        brand.footer = meta["title"]
    if not brand.version and meta.get("version"):
        brand.version = meta["version"]

    with tempfile.TemporaryDirectory(prefix="pptx_gen_") as tmp:
        tmp_dir = Path(tmp)

        logger.info("Phase 2: rendering mermaid diagrams")
        diagrams = render_all_diagrams(slides, tmp_dir)

        logger.info("Phase 3: assembling pptx")
        prs = load_template(template_path)

        for s in slides:
            if s.type == "architecture_diagram":
                build_diagram_slide(prs, s, diagrams)
            elif s.type == "kpi_slide":
                build_kpi_slide(prs, s, brand)
            else:
                builder = BUILDERS.get(s.type)
                if builder is None:
                    logger.error("slide %d: no builder for type %r", s.id, s.type)
                    continue
                builder(prs, s)
            logger.debug("  slide %d (%s) built", s.id, s.type)

        # Post-assembly: strip empty picture placeholders (prevents white boxes
        # covering template logos), then apply chrome + font uniformly.
        for slide in prs.slides:
            strip_unused_picture_placeholders(slide)
        apply_chrome(prs, slides, brand)
        apply_brand_font(prs, brand)

        saved = save_with_retry(prs, output_path)

    logger.info("Phase 4: verification")
    issues = verify_output(saved, expected_slides=len(slides))
    if issues:
        logger.warning("quality issues:")
        for msg in issues:
            logger.warning("  - %s", msg)
    else:
        logger.info("✓ all slides passed quality check")

    logger.info("✅ saved → %s (%d slides)", saved, len(slides))
    return saved


# ── CLI ───────────────────────────────────────────────────────────────
def _parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Generate PPTX from slides.json.")
    p.add_argument("--json", required=True, type=Path, help="path to slides.json")
    p.add_argument("--template", type=Path, default=None, help="optional .pptx template")
    p.add_argument(
        "--out",
        type=Path,
        default=Path("output_presentation.pptx"),
        help="output .pptx path (default: output_presentation.pptx)",
    )
    p.add_argument(
        "--brand-color", default=DEFAULT_BRAND_COLOR, help="brand color hex, e.g. #007A33"
    )
    p.add_argument("--font", default=None, help="override body/title font (e.g. 'Noto Sans TC')")
    p.add_argument("--footer", default=None, help="footer text (default: metadata.title)")
    p.add_argument("--version-label", default=None, help="version/date shown bottom-center")
    p.add_argument("--watermark", default=None, help="diagonal watermark text (e.g. CONFIDENTIAL)")
    p.add_argument("--page-numbers", action="store_true", help="show 'n / total' bottom-right")
    p.add_argument("-v", "--verbose", action="count", default=0, help="-v = INFO, -vv = DEBUG")
    return p.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = _parse_args(argv)
    level = (
        logging.WARNING
        if args.verbose == 0
        else logging.INFO if args.verbose == 1 else logging.DEBUG
    )
    logging.basicConfig(level=level, format="%(levelname)s %(name)s: %(message)s")

    brand = BrandConfig(
        color=args.brand_color,
        font=args.font,
        footer=args.footer,
        version=args.version_label,
        watermark=args.watermark,
        page_numbers=args.page_numbers,
    )

    try:
        generate(args.json, args.template, args.out, brand=brand)
    except FileNotFoundError as exc:
        logger.error("file not found: %s", exc)
        return 2
    except ValueError as exc:
        logger.error("invalid input: %s", exc)
        return 2
    return 0


if __name__ == "__main__":
    sys.exit(main())
